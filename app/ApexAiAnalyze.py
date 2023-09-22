import os
from pptx import Presentation
from pathlib import Path

from OpenAiChat import OpenAIChat


class ApexAiAnalyze:
    def __init__(self, assets_directory, prompt_path, destination_directory, configuration_directory):
        self.assets_directory = Path(assets_directory)
        self.prompt_path = prompt_path
        self.destination_directory = Path(destination_directory)
        self.configuration_directory = Path(configuration_directory)
        self.openai_chat = OpenAIChat(configuration_directory=self.configuration_directory)

    def _read_prompt(self):
        with open(self.prompt_path, 'r') as file:
            return file.read()

    def _generate_powerpoint(self, result):
        # Create a new presentation
        prs = Presentation()

        # Add a slide with title and content
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "OpenAI API Results"
        content.text = result

        # Save the presentation
        prs.save(os.path.join(self.destination_directory, "results.pptx"))

    def _call_openai_api(self, messages):
        return self.openai_chat.chat(messages)

    def process(self):
        # Read prompt from file
        with open(self.prompt_path, 'r') as file:
            prompt_content = file.read().strip()

        # Create a new PowerPoint presentation
        prs = Presentation()

        # Iterate through each file in the assets directory
        print("Assets directory ", self.assets_directory)
        for filepath in self.assets_directory.iterdir():
            if filepath.suffix in ['.csv', '.txt']:
                filename = filepath.name
                # For demonstration, let's assume each file contains some content we want to analyze.
                with open(os.path.join(self.assets_directory, filename), 'r') as file:
                    print(f"{filename} processing")
                    file_content = file.read().strip()

                    # Simulate sending content to the Chat completions API
                    messages = [
                        {"role": "system", "content": "You are a helpful assistant."},
                        {"role": "user", "content": prompt_content},
                        {"role": "user", "content": file_content}
                    ]
                    response = self._call_openai_api(messages)  # Call open AI with messages

                    # Add the response to the PowerPoint
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    title = slide.shapes.title
                    content = slide.placeholders[1]

                    title.text = f"Analysis for {filename}"
                    content.text = response['choices'][0]['message']['content']

        # Save the presentation
        prs.save(os.path.join(self.destination_directory, "results.pptx"))


def main():
    # Define your paths here
    assets_directory = "../data/jg23726"  # Example directory for assets
    prompt_path = "../data/prompts/prompt.txt"  # Example path for the prompt file
    destination_directory = "../data/results"  # Directory where the result PPT will be saved
    configuration_directory = "../data/configuration"  # Directory where the result PPT will be saved

    # Create directories if they don't exist
    os.makedirs(assets_directory, exist_ok=True)
    os.makedirs(destination_directory, exist_ok=True)

    # For testing purposes, let's create a dummy prompt.txt if it doesn't exist
    if not os.path.exists(prompt_path):
        with open(prompt_path, 'w') as f:
            f.write("This is a test prompt.")

    # Create an instance of the ApexAiAnalyze class and call process
    analyzer = ApexAiAnalyze(assets_directory, prompt_path, destination_directory, configuration_directory)
    analyzer.process()

    print(f"PowerPoint generated and saved in {destination_directory}/results.pptx")


if __name__ == "__main__":
    main()
